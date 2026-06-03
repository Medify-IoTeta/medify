from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colors ────────────────────────────────────────────────────────────────────
PRIMARY   = RGBColor(0x1A, 0x73, 0xE8)   # blue
DARK      = RGBColor(0x1E, 0x1E, 0x2E)   # almost black
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF4, 0xF7, 0xFF)
ACCENT    = RGBColor(0x34, 0xA8, 0x53)   # green (done items)
MUTED     = RGBColor(0x55, 0x65, 0x7A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=DARK,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def slide_header(slide, title, subtitle=None):
    """Dark top bar with title."""
    add_rect(slide, 0, 0, 13.33, 1.3, fill_color=DARK)
    add_text_box(slide, title,
                 left=0.4, top=0.15, width=10, height=0.7,
                 font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     left=0.4, top=0.82, width=12, height=0.4,
                     font_size=14, color=RGBColor(0xAA, 0xBB, 0xCC))


def bullet_block(slide, items, left, top, width=12.4, font_size=17,
                 color=DARK, marker="•", spacing=0.42):
    for i, item in enumerate(items):
        add_text_box(slide, f"{marker}  {item}",
                     left=left, top=top + i * spacing,
                     width=width, height=spacing + 0.05,
                     font_size=font_size, color=color)


def section_label(slide, text, left, top, width=3.5):
    box = add_rect(slide, left, top, width, 0.38, fill_color=PRIMARY)
    add_text_box(slide, text,
                 left=left + 0.1, top=top + 0.02,
                 width=width - 0.1, height=0.35,
                 font_size=13, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=DARK)
add_rect(slide, 0, 3.2, 13.33, 0.06, fill_color=PRIMARY)   # accent line

add_text_box(slide, "Medify",
             left=1, top=1.5, width=11, height=1.4,
             font_size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, "Smart Medication Management System",
             left=1, top=3.0, width=11, height=0.7,
             font_size=26, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)
add_text_box(slide, "IoT · Flutter · Spring Boot · Arduino",
             left=1, top=3.7, width=11, height=0.5,
             font_size=16, color=PRIMARY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Space
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_BG)
slide_header(slide, "The Problem", "Why medication management is a critical challenge")

problems = [
    "Patients taking multiple medications face high risk of missed or incorrect intake",
    "Passive solutions (pill organizers) provide no reminders or confirmation",
    "Caregivers and family members have no visibility into daily intake compliance",
    "Complex apps are not suitable for elderly or cognitively impaired users",
    "Non-compliance leads to serious medical consequences and hospitalizations",
]
bullet_block(slide, problems, left=0.6, top=1.55, font_size=18)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Our Solution
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_BG)
slide_header(slide, "Our Solution", "Medify — an end-to-end smart pill management system")

components = [
    ("Smart Pill Box",   "Automatically dispenses medication at scheduled times after user approval"),
    ("Mobile App",       "Reminds patients, confirms intake, and gives caregivers a monitoring dashboard"),
    ("Backend Server",   "Coordinates scheduling, intake tracking, notifications, and caregiver alerts"),
    ("Database",         "Stores full intake history, medication schedules, and user profiles"),
]

for i, (title, desc) in enumerate(components):
    top = 1.6 + i * 1.3
    add_rect(slide, 0.5, top, 12.3, 1.1, fill_color=WHITE, line_color=PRIMARY)
    add_text_box(slide, title,
                 left=0.75, top=top + 0.08, width=3.5, height=0.45,
                 font_size=16, bold=True, color=PRIMARY)
    add_text_box(slide, desc,
                 left=0.75, top=top + 0.52, width=11.5, height=0.5,
                 font_size=14, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Specific Features
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_BG)
slide_header(slide, "Key Features")

left_features = [
    "Three daily intake windows: Morning, Noon, Evening",
    "Real-time mobile reminders when an intake window opens",
    "Intake approval via app or physical device button",
    "Snooze (15 min or custom time) and Skip options",
]
right_features = [
    "Caregiver dashboard with real-time intake status",
    "Missed and incomplete intake alerts sent to caregivers",
    "Per-medication pre-intake reminders (e.g., 'Eat first')",
    "Full intake history — taken, missed, skipped, postponed",
]

section_label(slide, "Patient", left=0.5, top=1.45)
bullet_block(slide, left_features, left=0.6, top=1.92, width=5.9, font_size=16)

section_label(slide, "Caregiver & Monitoring", left=6.9, top=1.45)
bullet_block(slide, right_features, left=7.0, top=1.92, width=6.0, font_size=16)

add_rect(slide, 6.6, 1.4, 0.05, 5.8, fill_color=PRIMARY)   # vertical divider

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — General Architecture
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_BG)
slide_header(slide, "System Architecture")

boxes = [
    (0.4,  3.2, 2.8, 1.2, "Smart Pill Box",  "Arduino / C++",       PRIMARY),
    (4.5,  2.1, 3.3, 1.2, "Backend Server",  "Spring Boot / Java",  DARK),
    (4.5,  4.5, 3.3, 1.2, "PostgreSQL 15",   "Database",            MUTED),
    (9.8,  3.2, 3.0, 1.2, "Flutter App",     "Dart (iOS & Android)", PRIMARY),
]

for (l, t, w, h, title, sub, col) in boxes:
    add_rect(slide, l, t, w, h, fill_color=col)
    add_text_box(slide, title,
                 left=l+0.1, top=t+0.1, width=w-0.2, height=0.5,
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, sub,
                 left=l+0.1, top=t+0.62, width=w-0.2, height=0.4,
                 font_size=12, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

# Pill Box <-> Backend (horizontal arrow)
add_rect(slide, 3.2, 3.75, 1.3, 0.06, fill_color=MUTED)
add_text_box(slide, "WiFi / REST", left=3.2, top=3.55, width=1.3, height=0.3,
             font_size=11, color=MUTED, align=PP_ALIGN.CENTER)

# Backend <-> Flutter (horizontal arrow)
add_rect(slide, 7.8, 3.75, 2.0, 0.06, fill_color=MUTED)
add_text_box(slide, "REST API", left=7.8, top=3.55, width=2.0, height=0.3,
             font_size=11, color=MUTED, align=PP_ALIGN.CENTER)

# Backend <-> DB (vertical arrow)
add_rect(slide, 5.9, 3.3, 0.06, 1.2, fill_color=MUTED)
add_text_box(slide, "JPA", left=6.1, top=3.8, width=0.6, height=0.3,
             font_size=11, color=MUTED)

# Pill Box <-> Flutter (diagonal label — direct HTTP)
add_rect(slide, 3.3, 2.0, 6.6, 0.06, fill_color=ACCENT)
add_text_box(slide, "Local WiFi / HTTP — dispense command + confirmation",
             left=4.0, top=1.6, width=5.2, height=0.35,
             font_size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Technology Stack
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_BG)
slide_header(slide, "Technology Stack", "Choices and justifications")

stack = [
    ("Backend",   "Java 21 + Spring Boot 3.4",
     "Mature, enterprise-grade framework · Hexagonal architecture for clean separation of concerns · Strong REST & scheduling support"),
    ("Frontend",  "Dart + Flutter",
     "Single codebase → iOS & Android · Rich widget library · Fast development cycle · Suitable for accessible, patient-facing UIs"),
    ("Hardware",  "C++ + Arduino",
     "Industry standard for embedded / IoT · Extensive hardware libraries · Reliable motor and sensor control · Low-level WiFi (ESP)"),
    ("Database",  "PostgreSQL 15",
     "Relational model fits structured medical data · ACID compliance ensures data integrity · Schema managed by Flyway"),
]

for i, (layer, tech, reason) in enumerate(stack):
    top = 1.55 + i * 1.35
    add_rect(slide, 0.4, top, 2.2, 1.15, fill_color=PRIMARY)
    add_text_box(slide, layer,
                 left=0.45, top=top+0.1, width=2.1, height=0.4,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, tech,
                 left=0.45, top=top+0.55, width=2.1, height=0.45,
                 font_size=11, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

    add_rect(slide, 2.85, top, 10.0, 1.15, fill_color=WHITE, line_color=RGBColor(0xCC, 0xDD, 0xFF))
    add_text_box(slide, reason,
                 left=3.1, top=top+0.25, width=9.5, height=0.75,
                 font_size=14, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Domain: Smart Pill Box
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_BG)
slide_header(slide, "Domain — Smart Pill Box", "Arduino / C++")

features = [
    ("Medication Release",      "Motor-driven mechanism releases medication into the intake compartment at scheduled times, after user approval"),
    ("Compartment Monitoring",  "Sensor detects whether the intake compartment was emptied within the allowed time window"),
    ("Local Indication",        "LEDs and on-device display guide the user when it is time to take medication"),
    ("WiFi Communication",      "Communicates with the backend over WiFi using REST API to report intake status and receive updates"),
    ("Autonomous Operation",    "Continues functioning independently even when the patient's phone is not nearby"),
]

for i, (title, desc) in enumerate(features):
    top = 1.55 + i * 1.05
    add_rect(slide, 0.4, top, 3.0, 0.9, fill_color=PRIMARY)
    add_text_box(slide, title,
                 left=0.5, top=top+0.2, width=2.8, height=0.55,
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc,
                 left=3.6, top=top+0.18, width=9.4, height=0.6,
                 font_size=14, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Domain: Mobile App
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_BG)
slide_header(slide, "Domain — Mobile Application", "Flutter / Dart")

cols = [
    ("Patient View", [
        "Medication schedule grouped by window (Morning / Noon / Evening)",
        "Intake dialog with confirm, snooze (15 min or custom), and skip",
        "Progress ring showing completed windows",
        "Simple, accessible UI for elderly users",
    ]),
    ("Caregiver View", [
        "Real-time intake status for the monitored patient",
        "Summary card: windows completed today",
        "Expandable medication schedule",
        "Missed intake alerts from notification log",
    ]),
]

for ci, (title, items) in enumerate(cols):
    left = 0.5 + ci * 6.5
    section_label(slide, title, left=left, top=1.45, width=5.8)
    bullet_block(slide, items, left=left+0.1, top=1.95, width=5.7, font_size=15)

add_rect(slide, 6.6, 1.4, 0.05, 5.8, fill_color=PRIMARY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Domain: Backend
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_BG)
slide_header(slide, "Domain — Backend Server", "Java 21 · Spring Boot 3.4 · Hexagonal Architecture")

section_label(slide, "Module Structure", left=0.4, top=1.45, width=4.0)
structure = [
    "app  — Spring Boot entry point",
    "api  — REST controllers, services, NotificationAdapter",
    "domain  — Pure business logic, ports, schedulers, models",
    "data  — JPA repositories, Flyway migrations",
]
bullet_block(slide, structure, left=0.5, top=1.92, width=5.8, font_size=14, marker="▸")

section_label(slide, "Key Responsibilities", left=7.0, top=1.45, width=5.8)
responsibilities = [
    "Scheduled intake window creation (ReminderScheduler)",
    "Missed intake detection (MissedIntakeScheduler)",
    "Notification queue — polled by Flutter every 3 sec",
    "Intake status lifecycle: PENDING → APPROVED → TAKEN",
    "Caregiver alert logging (NotificationLog)",
    "REST API for app & device",
]
bullet_block(slide, responsibilities, left=7.1, top=1.92, width=5.9, font_size=14, marker="▸")

add_rect(slide, 6.6, 1.4, 0.05, 5.8, fill_color=PRIMARY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Domain: Database
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_BG)
slide_header(slide, "Domain — Database", "PostgreSQL 15 · Schema managed by Flyway")

tables = [
    ("users",            "User profiles, roles (PATIENT / CAREGIVER), authentication data"),
    ("medications",      "Medication name, timing window, dosage, pre-intake reminder rules"),
    ("intakes",          "One record per timing window per day — tracks approval, release, and status"),
    ("caregiver_links",  "Patient–caregiver associations and alert preferences"),
    ("notification_logs","History of all notifications and alerts sent by the system"),
]

for i, (table, desc) in enumerate(tables):
    top = 1.55 + i * 1.05
    add_rect(slide, 0.4, top, 3.2, 0.85, fill_color=DARK)
    add_text_box(slide, table,
                 left=0.5, top=top+0.2, width=3.0, height=0.5,
                 font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc,
                 left=3.85, top=top+0.22, width=9.1, height=0.5,
                 font_size=15, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Current Status
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_BG)
slide_header(slide, "Current Project Status")

section_label(slide, "Completed", left=0.4, top=1.45, width=3.0)
done = [
    "Backend core — medications, intakes, notifications, caregiver links",
    "Automatic scheduling — Morning 08:00, Noon 13:00, Evening 20:00",
    "Missed intake detection & alert logging",
    "Flutter patient screen — schedule, intake dialog, snooze/skip",
    "Flutter caregiver screen — intake status, missed alerts",
    "Arduino device — WiFi, motor release, physical button",
    "PostgreSQL with Flyway-managed schema",
]
for i, item in enumerate(done):
    add_text_box(slide, f"✓  {item}",
                 left=0.5, top=1.92 + i*0.62, width=12.3, height=0.55,
                 font_size=15, color=ACCENT)

section_label(slide, "Planned", left=0.4, top=6.22, width=3.0)
planned = [
    "User authentication  ·  Intake compartment sensor (INCOMPLETE detection)  ·  Push notifications  ·  Device management API",
]
add_text_box(slide, f"○  {planned[0]}",
             left=0.5, top=6.55, width=12.3, height=0.5,
             font_size=14, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — AI Usage
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_BG)
slide_header(slide, "AI Usage in Development", "Where we used AI and what we added ourselves")

section_label(slide, "AI-Assisted", left=0.4, top=1.45, width=3.5)
ai_items = [
    "Code generation — boilerplate, domain models, REST controllers",
    "Architecture review — hexagonal structure, port/adapter design",
    "Spec alignment — comparing implementation against project overview",
    "Debugging — diagnosing runtime and integration issues",
    "Documentation — CLAUDE.md files, README",
]
bullet_block(slide, ai_items, left=0.5, top=1.92, width=5.8, font_size=15)

section_label(slide, "Our Contribution", left=7.0, top=1.45, width=5.8)
our_items = [
    "System design decisions — IoT architecture, hardware selection",
    "Domain logic — intake lifecycle rules, timing windows",
    "Hardware implementation — Arduino wiring, motor control, sensor",
    "UI/UX decisions — patient-accessible design, caregiver flow",
    "Integration — end-to-end testing, WiFi communication",
]
bullet_block(slide, our_items, left=7.1, top=1.92, width=5.9, font_size=15)

add_rect(slide, 6.6, 1.4, 0.05, 5.8, fill_color=PRIMARY)

# ── Save ──────────────────────────────────────────────────────────────────────
output_path = "C:/Users/netta/medify/docs/Medify_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
